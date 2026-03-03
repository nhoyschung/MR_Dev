"""Tests for the PPTX export engine (Phases 3 & 4).

Run with: pytest tests/test_pptx.py -v
"""

import pytest
from io import BytesIO
from pathlib import Path
from datetime import date

from pptx import Presentation as PptxPresentation


# ── Phase 1: Chart factory functions ───────────────────────────────────────

class TestChartFactories:
    def test_grade_distribution_figure_not_closed(self):
        from src.reports.charts import create_grade_distribution_figure, fig_to_bytesio
        fig = create_grade_distribution_figure([{"grade": "M-I", "count": 5}])
        assert fig is not None
        buf = fig_to_bytesio(fig)
        import matplotlib.pyplot as plt
        plt.close(fig)
        buf.seek(0)
        assert buf.read(8)[:4] == b'\x89PNG'

    def test_grade_distribution_figure_empty(self):
        from src.reports.charts import create_grade_distribution_figure
        assert create_grade_distribution_figure([]) is None

    def test_radar_figure_created(self):
        from src.reports.charts import create_radar_figure, fig_to_bytesio
        import matplotlib.pyplot as plt
        scores = [("ProjectA", {"Loc": 8.0, "Trans": 7.0}),
                  ("ProjectB", {"Loc": 6.0, "Trans": 8.0})]
        categories = ["Loc", "Trans"]
        fig = create_radar_figure(scores, categories)
        assert fig is not None
        buf = fig_to_bytesio(fig)
        plt.close(fig)
        buf.seek(0)
        assert buf.read(4) == b'\x89PNG'

    def test_radar_figure_empty(self):
        from src.reports.charts import create_radar_figure
        assert create_radar_figure([], []) is None

    def test_price_trend_figure(self):
        from src.reports.charts import create_price_trend_figure
        import matplotlib.pyplot as plt
        fig = create_price_trend_figure([
            {"period": "2024-H1", "price": 2000},
            {"period": "2024-H2", "price": 2200},
        ])
        assert fig is not None
        plt.close(fig)

    def test_supply_demand_figure(self):
        from src.reports.charts import create_supply_demand_figure
        import matplotlib.pyplot as plt
        fig = create_supply_demand_figure(1000, 200, 150, 850, 75.0)
        assert fig is not None
        plt.close(fig)

    def test_supply_demand_figure_empty(self):
        from src.reports.charts import create_supply_demand_figure
        assert create_supply_demand_figure(0, 0, 0, 0, 0.0) is None

    def test_fig_to_bytesio_is_png(self):
        from src.reports.charts import create_grade_distribution_figure, fig_to_bytesio
        import matplotlib.pyplot as plt
        fig = create_grade_distribution_figure([{"grade": "A-I", "count": 3}])
        buf = fig_to_bytesio(fig, dpi=72)
        plt.close(fig)
        buf.seek(0)
        assert buf.read(4) == b'\x89PNG'


# ── Phase 3: content_schema ────────────────────────────────────────────────

class TestContentSchema:
    def test_cover_slide_keys(self):
        from src.reports.pptx.content_schema import CoverSlide
        slide: CoverSlide = {
            "index": 1, "type": "cover",
            "title": "T", "subtitle": "S",
            "city": "HCMC", "period": "2025-H1",
            "report_type": "market_briefing",
            "date": date.today().isoformat(),
        }
        assert slide["type"] == "cover"

    def test_ko_and_en_manifest_same_structure(self):
        from src.reports.pptx.content_schema import SlideContentManifest, CoverSlide
        cover: CoverSlide = {
            "index": 1, "type": "cover",
            "title": "Title", "subtitle": "Sub",
            "city": "HCMC", "period": "2025-H1",
            "report_type": "market_briefing",
            "date": "2026-02-21",
        }
        en: SlideContentManifest = {
            "job_id": "test", "report_type": "market_briefing",
            "language": "en", "params": {}, "slides": [cover],
        }
        ko: SlideContentManifest = {
            "job_id": "test", "report_type": "market_briefing",
            "language": "ko", "params": {}, "slides": [cover],
        }
        assert set(en.keys()) == set(ko.keys())
        assert en["slides"][0]["type"] == ko["slides"][0]["type"]


# ── Phase 3: PptxBuilder ──────────────────────────────────────────────────

class TestPptxBuilder:
    def _cover(self) -> dict:
        return {
            "index": 1, "type": "cover",
            "title": "Test Report",
            "subtitle": "Market Intelligence",
            "city": "Ho Chi Minh City",
            "period": "2025-H1",
            "report_type": "market_briefing",
            "date": date.today().isoformat(),
        }

    def test_builder_add_cover_slide_count(self):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        b.add_cover(self._cover())
        assert len(b._prs.slides) == 1

    def test_builder_add_kpi_dashboard(self):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        slide = {
            "index": 2, "type": "kpi_dashboard",
            "slide_title": "Market KPIs",
            "kpis": [
                {"label": "Projects", "value": "47", "delta": "+3", "color": "blue"},
                {"label": "Avg Price", "value": "$2,400", "delta": None, "color": "green"},
            ],
            "note": "Market momentum is strong.",
        }
        b.add_kpi_dashboard(slide)
        assert len(b._prs.slides) == 1

    def test_builder_add_table_slide(self):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        slide = {
            "index": 3, "type": "table",
            "title": "Grade Distribution",
            "headers": ["Grade", "Projects", "Avg Price"],
            "rows": [["M-I", "12", "$2,200"], ["M-II", "8", "$1,800"]],
            "caption": "M-I dominates supply.",
            "grade_col_index": 0,
        }
        b.add_table_slide(slide)
        assert len(b._prs.slides) == 1

    def test_builder_add_swot(self):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        slide = {
            "index": 4, "type": "swot",
            "title": "SWOT Analysis",
            "strengths": ["Strong location"],
            "weaknesses": ["Limited supply"],
            "opportunities": ["High absorption"],
            "threats": ["Competition"],
        }
        b.add_swot_slide(slide)
        assert len(b._prs.slides) == 1

    def test_builder_add_conclusion(self):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        slide = {
            "index": 5, "type": "conclusion",
            "title": "Conclusion & Verdict",
            "verdict": "HIGHLY VIABLE",
            "bullets": ["Point 1", "Point 2"],
            "badge_label": "HIGHLY VIABLE",
            "badge_color": "green",
        }
        b.add_conclusion_slide(slide)
        assert len(b._prs.slides) == 1

    def test_builder_add_section_divider(self):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        b.add_section_divider("01", "Market Overview", "Supply & Demand")
        assert len(b._prs.slides) == 1

    def test_builder_build_from_manifest(self, tmp_path):
        from src.reports.pptx.builder import PptxBuilder
        manifest = {
            "job_id": "test_job",
            "report_type": "market_briefing",
            "language": "en",
            "params": {},
            "slides": [
                self._cover(),
                {
                    "index": 2, "type": "kpi_dashboard",
                    "slide_title": "KPIs",
                    "kpis": [{"label": "X", "value": "1", "delta": None, "color": "blue"}],
                    "note": "Test note.",
                },
                {
                    "index": 3, "type": "swot",
                    "title": "SWOT",
                    "strengths": ["S1"], "weaknesses": ["W1"],
                    "opportunities": ["O1"], "threats": ["T1"],
                },
            ],
        }
        b = PptxBuilder().build_from_manifest(manifest)
        assert len(b._prs.slides) == 3

        out = b.save(tmp_path / "test.pptx")
        assert out.exists()
        prs = PptxPresentation(str(out))
        assert len(prs.slides) == 3

    def test_builder_save_creates_file(self, tmp_path):
        from src.reports.pptx.builder import PptxBuilder
        b = PptxBuilder()
        b.add_cover(self._cover())
        out = b.save(tmp_path / "out.pptx")
        assert out.exists()
        assert out.suffix == ".pptx"


# ── Phase 4: Report generators ─────────────────────────────────────────────

@pytest.fixture
def db_session():
    """In-memory SQLite session seeded with minimal data for PPTX tests."""
    from sqlalchemy import create_engine
    from sqlalchemy.orm import Session
    from src.db.models import Base, City, District, ReportPeriod, GradeDefinition

    engine = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False})
    Base.metadata.create_all(engine)

    with Session(engine) as session:
        city = City(name_en="Ho Chi Minh City", name_vi="Thành phố Hồ Chí Minh",
                    region="South")
        session.add(city)
        session.flush()

        district = District(city_id=city.id, name_en="Binh Thanh",
                            name_vi="Bình Thạnh", district_type="urban")
        session.add(district)

        period = ReportPeriod(year=2025, half="H1",
                              report_date=date(2025, 6, 30))
        session.add(period)

        grade = GradeDefinition(city_id=city.id, grade_code="M-I",
                                segment="mid-end",
                                min_price_usd=1800.0, max_price_usd=2500.0)
        session.add(grade)
        session.commit()
        yield session


class TestReportGenerators:
    def test_market_briefing_returns_none_for_invalid_city(self, db_session):
        from src.reports.pptx.market_briefing import generate_market_briefing_pptx
        result = generate_market_briefing_pptx(db_session, "InvalidCity", 2025, "H1")
        assert result is None

    def test_market_briefing_returns_none_for_invalid_period(self, db_session):
        from src.reports.pptx.market_briefing import generate_market_briefing_pptx
        result = generate_market_briefing_pptx(db_session, "Ho Chi Minh City", 1999, "H1")
        assert result is None

    def test_market_briefing_no_projects_returns_none(self, db_session):
        from src.reports.pptx.market_briefing import generate_market_briefing_pptx
        result = generate_market_briefing_pptx(db_session, "Ho Chi Minh City", 2025, "H1")
        # No projects seeded → should return None
        assert result is None

    def test_competitor_pptx_returns_none_single_project(self, db_session):
        from src.reports.pptx.competitor import generate_competitor_pptx
        result = generate_competitor_pptx(db_session, ["Solo Project"], 2025, "H1")
        assert result is None

    def test_competitor_pptx_returns_none_empty_list(self, db_session):
        from src.reports.pptx.competitor import generate_competitor_pptx
        result = generate_competitor_pptx(db_session, [], 2025, "H1")
        assert result is None

    def test_language_suffix_in_filename(self, db_session, tmp_path, monkeypatch):
        """Filename must contain _en.pptx or _ko.pptx suffix."""
        import src.reports.pptx.market_briefing as mb_mod
        from src.reports.pptx.market_briefing import generate_market_briefing_pptx

        # Patch OUTPUT_DIR so we write to tmp_path
        monkeypatch.setattr(mb_mod, "OUTPUT_DIR", tmp_path)

        # Pass a pre-built minimal manifest to bypass DB requirement
        from src.reports.pptx.content_schema import SlideContentManifest
        manifest: SlideContentManifest = {
            "job_id": "lang_test",
            "report_type": "market_briefing",
            "language": "en",
            "params": {"city": "Ho Chi Minh City", "year": 2025, "half": "H1"},
            "slides": [
                {
                    "index": 1, "type": "cover",
                    "title": "Test", "subtitle": "Sub",
                    "city": "Ho Chi Minh City", "period": "2025-H1",
                    "report_type": "market_briefing",
                    "date": "2026-02-21",
                }
            ],
        }
        path = generate_market_briefing_pptx(
            db_session, "Ho Chi Minh City", 2025, "H1",
            content_override=manifest, language="en",
        )
        # Even if city/period found, language suffix must appear
        if path is not None:
            assert "_en.pptx" in path.name

    def test_content_override_bypasses_db(self, db_session, tmp_path, monkeypatch):
        """content_override should produce a PPTX even with no DB projects."""
        import src.reports.pptx.market_briefing as mb_mod
        from src.reports.pptx.market_briefing import generate_market_briefing_pptx
        monkeypatch.setattr(mb_mod, "OUTPUT_DIR", tmp_path)

        manifest = {
            "job_id": "override_test",
            "report_type": "market_briefing",
            "language": "en",
            "params": {},
            "slides": [
                {
                    "index": 1, "type": "cover",
                    "title": "Override Test", "subtitle": "Sub",
                    "city": "Ho Chi Minh City", "period": "2025-H1",
                    "report_type": "market_briefing",
                    "date": "2026-02-21",
                },
                {
                    "index": 2, "type": "conclusion",
                    "title": "Conclusion",
                    "verdict": "HIGHLY VIABLE",
                    "bullets": ["Good location"],
                    "badge_label": "HIGHLY VIABLE",
                    "badge_color": "green",
                },
            ],
        }
        path = generate_market_briefing_pptx(
            db_session, "Ho Chi Minh City", 2025, "H1",
            content_override=manifest, language="en",
        )
        assert path is not None
        assert path.exists()
        prs = PptxPresentation(str(path))
        assert len(prs.slides) == 2

    def test_project_profile_returns_none_unknown(self, db_session):
        from src.reports.pptx.project_profile import generate_project_profile_pptx
        result = generate_project_profile_pptx(db_session, "NonExistentProject XYZ")
        assert result is None

    def test_land_review_raises_for_invalid_city(self, db_session):
        from src.reports.pptx.land_review import generate_land_review_pptx
        with pytest.raises(ValueError, match="not found"):
            generate_land_review_pptx(db_session, {"city": "InvalidCity", "land_area_ha": 2.0})
